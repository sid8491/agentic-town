"""Build a presentation deck explaining the Gurgaon Town Life codebase.

Usage:
    .venv/Scripts/python.exe scripts/build_ppt.py
    -> writes docs/gurgaon_town_life_explained.pptx
"""

from __future__ import annotations

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------

BG       = RGBColor(0x10, 0x14, 0x1F)   # near-black navy
PANEL    = RGBColor(0x18, 0x20, 0x32)   # card background
INK      = RGBColor(0xE8, 0xEC, 0xF4)   # primary text
DIM      = RGBColor(0x9A, 0xA4, 0xB8)   # secondary text
ACCENT   = RGBColor(0xFF, 0x9E, 0x3D)   # warm Gurgaon orange
ACCENT2  = RGBColor(0x4D, 0xA8, 0xDA)   # cool blue
GOOD     = RGBColor(0x5D, 0xD3, 0x9C)   # mint green
LINE     = RGBColor(0x2A, 0x35, 0x4A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def add_bg(slide, color=BG):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.shadow.inherit = False
    return rect


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        lines = [text]
    else:
        lines = list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_panel(slide, left, top, width, height, fill=PANEL, line=LINE):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.05
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = line
    rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    return rect


def add_chip(slide, text, left, top, width, height, fill=ACCENT, color=BG, size=12, bold=True):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.4
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.fill.background()
    rect.shadow.inherit = False
    tf = rect.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return rect


def slide_header(slide, kicker, title):
    """Standard kicker + title block at the top of a content slide."""
    add_text(slide, kicker.upper(), Inches(0.6), Inches(0.35),
             Inches(12), Inches(0.3), size=12, bold=True, color=ACCENT)
    add_text(slide, title, Inches(0.6), Inches(0.65),
             Inches(12), Inches(0.7), size=28, bold=True, color=INK)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), Inches(1.35),
                                  Inches(1.2), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    line.shadow.inherit = False


def add_bullets(slide, items, left, top, width, height,
                size=16, color=INK, bullet_color=ACCENT, line_height=1.2):
    """Each item is a string OR a (head, sub) tuple."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.05)
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            head, sub = item
        else:
            head, sub = item, None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        bullet = p.add_run()
        bullet.text = "▸ "
        bullet.font.size = Pt(size)
        bullet.font.bold = True
        bullet.font.color.rgb = bullet_color
        bullet.font.name = "Calibri"
        run = p.add_run()
        run.text = head
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = color
        if sub:
            sp = tf.add_paragraph()
            sp.alignment = PP_ALIGN.LEFT
            sp.space_after = Pt(6)
            sub_run = sp.add_run()
            sub_run.text = "    " + sub
            sub_run.font.name = "Calibri"
            sub_run.font.size = Pt(size - 2)
            sub_run.font.color.rgb = DIM
    return box


def add_code(slide, code, left, top, width, height, size=12):
    panel = add_panel(slide, left, top, width, height,
                      fill=RGBColor(0x0C, 0x10, 0x18), line=LINE)
    box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1),
                                   width - Inches(0.3), height - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    lines = code.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = "Consolas"
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(0xCD, 0xE3, 0xFF)
    return panel


def add_footer(slide, page_num, total):
    add_text(slide, "Gurgaon Town Life — Codebase Walkthrough",
             Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             size=10, color=DIM)
    add_text(slide, f"{page_num} / {total}",
             Inches(11.5), Inches(7.05), Inches(1.2), Inches(0.3),
             size=10, color=DIM, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

slides_built = []


def new_slide():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    slides_built.append(s)
    return s


# ---- 1. Title ------------------------------------------------------------

s = new_slide()
# Big accent bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8),
                         prs.slide_width, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background(); bar.shadow.inherit = False

add_text(s, "GURGAON TOWN LIFE",
         Inches(0.8), Inches(1.6), Inches(12), Inches(1.0),
         size=54, bold=True, color=INK)
add_text(s, "An autonomous multi-agent simulation — explained from the ground up",
         Inches(0.8), Inches(3.0), Inches(12), Inches(0.6),
         size=22, color=DIM)
add_text(s, "10 LLM-powered characters · LangGraph · FastAPI · Arcade · Ollama / Gemini",
         Inches(0.8), Inches(3.7), Inches(12), Inches(0.5),
         size=16, color=ACCENT2)

# Three chips
add_chip(s, "PYTHON 3.12", Inches(0.8), Inches(5.0), Inches(1.6), Inches(0.4),
         fill=PANEL, color=INK, size=12)
add_chip(s, "LANGGRAPH", Inches(2.5), Inches(5.0), Inches(1.6), Inches(0.4),
         fill=PANEL, color=INK, size=12)
add_chip(s, "ASYNCIO", Inches(4.2), Inches(5.0), Inches(1.6), Inches(0.4),
         fill=PANEL, color=INK, size=12)
add_chip(s, "FASTAPI", Inches(5.9), Inches(5.0), Inches(1.6), Inches(0.4),
         fill=PANEL, color=INK, size=12)
add_chip(s, "ARCADE", Inches(7.6), Inches(5.0), Inches(1.6), Inches(0.4),
         fill=PANEL, color=INK, size=12)
add_chip(s, "LITELLM", Inches(9.3), Inches(5.0), Inches(1.6), Inches(0.4),
         fill=PANEL, color=INK, size=12)

add_text(s, "A walkthrough of the codebase + how to present it in an interview",
         Inches(0.8), Inches(6.6), Inches(12), Inches(0.4),
         size=14, color=DIM)


# ---- 2. What is this? ----------------------------------------------------

s = new_slide()
slide_header(s, "Slide 1", "What is this project?")

add_text(s,
         "A virtual town where 10 AI characters live their own lives. "
         "Nobody plays them. You just watch.",
         Inches(0.6), Inches(1.6), Inches(12), Inches(1.0),
         size=22, color=INK)

# Three columns
col_w = Inches(4.0); col_h = Inches(3.6); top = Inches(2.9)
gap = Inches(0.2); left = Inches(0.6)

def feature_card(left, title, lines, accent=ACCENT):
    add_panel(s, left, top, col_w, col_h)
    add_text(s, title, left + Inches(0.3), top + Inches(0.25),
             col_w - Inches(0.6), Inches(0.5),
             size=18, bold=True, color=accent)
    add_text(s, lines, left + Inches(0.3), top + Inches(0.85),
             col_w - Inches(0.6), col_h - Inches(1.0),
             size=14, color=INK)

feature_card(Inches(0.6),
             "10 AI agents",
             ["Each has a unique personality file (soul.md).",
              "Each has private memory, diary, and goals — written by themselves.",
              "Backed by a local LLM (Ollama) or Gemini."])
feature_card(Inches(4.8),
             "A real little town",
             ["8 connected locations: apartment, metro, Cyber City, dhaba, market…",
              "Hunger, energy, mood decay over time.",
              "Money, items, conversations — all persisted."],
             accent=ACCENT2)
feature_card(Inches(9.0),
             "Two views",
             ["Arcade desktop window (Pythonic 2D renderer).",
              "Browser viewer at localhost:8000 — vanilla JS canvas.",
              "Both are read-only consumers of state."],
             accent=GOOD)


# ---- 3. The Big Picture (architecture) -----------------------------------

s = new_slide()
slide_header(s, "Slide 2", "Three strictly-separated layers")

# Layer boxes
layer_w = Inches(4.0); layer_h = Inches(2.4); ltop = Inches(1.8)

def layer(left, color, name, role, files):
    add_panel(s, left, ltop, layer_w, layer_h)
    chip_w = Inches(1.6)
    add_chip(s, name, left + Inches(0.3), ltop + Inches(0.25),
             chip_w, Inches(0.4), fill=color, color=BG, size=12)
    add_text(s, role, left + Inches(0.3), ltop + Inches(0.85),
             layer_w - Inches(0.6), Inches(0.6),
             size=14, color=INK, bold=True)
    add_text(s, files, left + Inches(0.3), ltop + Inches(1.4),
             layer_w - Inches(0.6), Inches(1.0),
             size=12, color=DIM)

layer(Inches(0.6), ACCENT, "ENGINE",
      "All simulation logic.\nNo rendering. No HTTP.",
      "engine/world.py · agent.py · tools.py\nllm.py · needs.py · narrator.py …")

layer(Inches(4.8), ACCENT2, "RENDERER",
      "Arcade window.\nReads state every frame.",
      "main.py — sprites, camera, UI panels.\nKey/click → engine, never directly writes state.")

layer(Inches(9.0), GOOD, "WEB",
      "FastAPI + vanilla JS.\nRuns in a background thread.",
      "server.py exposes /api/state etc.\nviewer.html draws on HTML5 canvas.")

# Arrow row
add_text(s, "All three look at the SAME WorldState object — single source of truth",
         Inches(0.6), Inches(4.5), Inches(12), Inches(0.5),
         size=16, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)

# State box
add_panel(s, Inches(3.5), Inches(5.2), Inches(6.3), Inches(1.6), fill=PANEL)
add_text(s, "WorldState  ➜  world/state.json",
         Inches(3.5), Inches(5.3), Inches(6.3), Inches(0.5),
         size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
add_text(s, "Sim time · agent positions · needs · inventories · inboxes · "
            "events log · plans · relationships",
         Inches(3.7), Inches(5.85), Inches(5.9), Inches(0.9),
         size=12, color=DIM, align=PP_ALIGN.CENTER)


# ---- 4. The 10 Characters -----------------------------------------------

s = new_slide()
slide_header(s, "Slide 3", "Meet the 10 residents")

people = [
    ("Arjun",  "Software engineer, 28",       "Startup, anxious, overworked"),
    ("Priya",  "Product manager, 32",         "MNC, organised, networker"),
    ("Rahul",  "Delivery boy, 22",            "Zomato, street-smart"),
    ("Kavya",  "Freelance designer, 26",      "Creative, odd hours"),
    ("Suresh", "Cab driver, 45",              "Wise, knows everyone"),
    ("Neha",   "HR professional, 30",         "Cheerful, gossip queen"),
    ("Vikram", "Retired colonel, 62",         "Disciplined, opinionated"),
    ("Deepa",  "Homemaker, 38",               "Resourceful, anchor"),
    ("Rohan",  "MDI student, 24",             "Idealistic, broke"),
    ("Anita",  "Boutique owner, 41",          "Entrepreneurial, proud"),
]

# 5 cols x 2 rows
card_w = Inches(2.4); card_h = Inches(2.2)
start_l = Inches(0.6); start_t = Inches(1.7)
gap_x = Inches(0.1); gap_y = Inches(0.2)
for i, (name, who, tag) in enumerate(people):
    col = i % 5; row = i // 5
    left = start_l + col * (card_w + gap_x)
    top  = start_t + row * (card_h + gap_y)
    add_panel(s, left, top, card_w, card_h)
    add_text(s, name, left + Inches(0.2), top + Inches(0.15),
             card_w - Inches(0.4), Inches(0.5),
             size=20, bold=True, color=ACCENT)
    add_text(s, who, left + Inches(0.2), top + Inches(0.7),
             card_w - Inches(0.4), Inches(0.5),
             size=12, color=INK)
    add_text(s, tag, left + Inches(0.2), top + Inches(1.2),
             card_w - Inches(0.4), Inches(0.9),
             size=11, color=DIM)

add_text(s, "Each character lives in agents/{name}/ — four small markdown files define their entire mind.",
         Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         size=13, color=DIM, align=PP_ALIGN.CENTER)


# ---- 5. Each agent's "mind" = 4 files ------------------------------------

s = new_slide()
slide_header(s, "Slide 4", "Each agent's mind is just four markdown files")

files = [
    ("soul.md",   "Personality, voice, values.",            "Written by developer.\nNEVER changes at runtime.", ACCENT),
    ("memory.md", "Long-term beliefs and relationships.",   "The agent rewrites this when something significant happens.", ACCENT2),
    ("diary.md",  "Daily journal — append-only.",            "The agent appends a reflection at end-of-day, or after big events.", GOOD),
    ("goals.md",  "Current priorities for today.",           "The agent rewrites this each morning.", RGBColor(0xE6, 0x6E, 0xC9)),
]

card_w = Inches(2.95); card_h = Inches(4.6); top = Inches(1.7)
for i, (name, what, who, color) in enumerate(files):
    left = Inches(0.6) + i * (card_w + Inches(0.1))
    add_panel(s, left, top, card_w, card_h)
    # Header strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, card_w, Inches(0.5))
    strip.fill.solid(); strip.fill.fore_color.rgb = color
    strip.line.fill.background(); strip.shadow.inherit = False
    add_text(s, name, left, top, card_w, Inches(0.5),
             size=18, bold=True, color=BG, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
    add_text(s, what, left + Inches(0.2), top + Inches(0.7),
             card_w - Inches(0.4), Inches(1.4),
             size=14, color=INK, bold=True)
    add_text(s, who, left + Inches(0.2), top + Inches(2.0),
             card_w - Inches(0.4), Inches(2.5),
             size=12, color=DIM)

add_text(s, "Implication: an agent's identity is GIT-versioned text. "
            "You can read, diff, and edit any character's mind by hand.",
         Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
         size=13, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)


# ---- 6. The 4-node decision loop ----------------------------------------

s = new_slide()
slide_header(s, "Slide 5", "The decision loop — every tick, every agent")

nodes = [
    ("gather_context",  "Read soul, memory excerpt,\nneeds, what's nearby, time", ACCENT),
    ("llm_decide",      "Send context + tool schemas\nto Ollama / Gemini", ACCENT2),
    ("execute_tool",    "Run the chosen tool\n(move, talk, eat, work…)", GOOD),
    ("reflect",         "Update last_action.\nMaybe append diary / memory.", RGBColor(0xE6, 0x6E, 0xC9)),
]

box_w = Inches(2.7); box_h = Inches(2.0); ntop = Inches(2.0)
for i, (name, body, color) in enumerate(nodes):
    left = Inches(0.6) + i * (box_w + Inches(0.4))
    add_panel(s, left, ntop, box_w, box_h)
    add_text(s, name, left, ntop + Inches(0.2), box_w, Inches(0.5),
             size=18, bold=True, color=color, align=PP_ALIGN.CENTER, font="Consolas")
    add_text(s, body, left + Inches(0.2), ntop + Inches(0.85),
             box_w - Inches(0.4), Inches(1.1),
             size=13, color=INK, align=PP_ALIGN.CENTER)
    if i < 3:
        # arrow
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 left + box_w + Inches(0.05),
                                 ntop + Inches(0.85),
                                 Inches(0.3), Inches(0.3))
        arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT
        arr.line.fill.background(); arr.shadow.inherit = False

add_text(s, "All ten agents run this graph in parallel each tick via asyncio.gather.",
         Inches(0.6), Inches(4.4), Inches(12), Inches(0.5),
         size=15, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)

add_text(s, "Why LangGraph? It gives a clean, composable state machine — "
            "each node is a pure function over a typed state dict. "
            "Easy to add nodes (e.g. a 'plan' node) without touching the others.",
         Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.0),
         size=14, color=DIM, align=PP_ALIGN.CENTER)

add_code(s,
         "graph = StateGraph(AgentState)\n"
         "graph.add_node('gather_context', gather_context)\n"
         "graph.add_node('llm_decide',     llm_decide)\n"
         "graph.add_node('execute_tool',   execute_tool)\n"
         "graph.add_node('reflect',        reflect)\n"
         "graph.set_entry_point('gather_context')\n"
         "graph.add_edge('gather_context', 'llm_decide')\n"
         "graph.add_edge('llm_decide',     'execute_tool')\n"
         "graph.add_edge('execute_tool',   'reflect')\n"
         "graph.add_edge('reflect',        END)",
         Inches(2.0), Inches(6.0), Inches(9.3), Inches(1.3), size=11)


# ---- 7. WorldState ------------------------------------------------------

s = new_slide()
slide_header(s, "Slide 6", "engine/world.py — the only place state lives")

add_bullets(s, [
    ("Single source of truth",
     "All agent positions, needs, inventories, inboxes, sim time, plans, "
     "relationships, events live in one dict, persisted to world/state.json."),
    ("Thread-safe via asyncio.Lock",
     "Every mutation acquires self._lock — ten agents writing concurrently never corrupt state."),
    ("Append-only event log",
     "Every action becomes an event row. Lets the renderer animate, the narrator summarise, "
     "and the drama-pacer score recent activity."),
    ("Save cadence",
     "After each full tick, world.save() atomically writes state.json (write to temp + rename)."),
    ("Reset semantics",
     "python main.py --reset wipes state.json but leaves agents/*/soul.md alone — "
     "factory-fresh world, same personalities."),
], Inches(0.6), Inches(1.7), Inches(7.6), Inches(5.5), size=14)

# Right side — tiny code sample
add_panel(s, Inches(8.5), Inches(1.7), Inches(4.2), Inches(5.0))
add_text(s, "Mutation pattern", Inches(8.7), Inches(1.85),
         Inches(3.8), Inches(0.4), size=14, bold=True, color=ACCENT)
add_code(s,
         "async with world._lock:\n"
         "  agent = world._state\n"
         "      ['agents'][name]\n"
         "  agent['hunger'] -= 30\n"
         "  agent['inventory']\n"
         "      .remove('bread')\n"
         "  world._state['events']\n"
         "      .append({...})\n"
         "await world.save()",
         Inches(8.7), Inches(2.3), Inches(3.8), Inches(4.2), size=12)


# ---- 8. agent.py details --------------------------------------------------

s = new_slide()
slide_header(s, "Slide 7", "engine/agent.py — what makes a character tick")

add_bullets(s, [
    ("AgentRunner class",
     "One per character. Holds the agent's name and the compiled LangGraph."),
    ("Per-archetype schedules",
     "office_worker, vendor, retired, homemaker, entrepreneur, student, night_owl. "
     "A SCHEDULE block is injected into the prompt based on sim_time — "
     "no extra LLM call needed to know 'I should sleep now'."),
    ("Sleep, work, lunch, evening windows",
     "Hard-coded ranges per archetype. Agents can override only if hunger/energy hit critical."),
    ("Workplace mapping",
     "Arjun goes to Cyber City, Anita to Cyber Hub, Suresh drives between "
     "apartment / metro / Sector 29 — never to Cyber City."),
    ("Memory grep",
     "Before deciding, the agent grep_memory()s for keywords from the current context "
     "(nearby people, location) — keeps prompts small even as memory.md grows."),
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), size=14)


# ---- 9. tools.py — the action catalog -----------------------------------

s = new_slide()
slide_header(s, "Slide 8", "engine/tools.py — every action an agent can take")

tool_groups = [
    ("Movement & perception",
     ["move_to(location)  — full BFS pathfinding, walks all hops in one call",
      "look_around()      — reachable locations + who is here",
      "check_needs()      — hunger / energy / mood",
      "check_inventory()  — items + cash"]),
    ("Social",
     ["talk_to(target, message)",
      "ask_about(target, topic)",
      "refuse / disagree(target, …)",
      "propose_plan / confirm_plan / decline_plan"]),
    ("Economy & survival",
     ["buy / sell(item, qty[, price])",
      "give_item(target, item, qty)",
      "eat(item)        — restores hunger from inventory",
      "eat_out()        — restaurant meal, costs cash",
      "work()           — earn money at workplace",
      "sleep_action()   — restore energy"]),
    ("Self-reflection",
     ["read_file(file)              — soul / memory / diary / goals",
      "edit_file(file, content)     — only memory.md / goals.md",
      "append_diary(entry)          — daily reflection",
      "grep_memory(query)"]),
]

for i, (group, items) in enumerate(tool_groups):
    col = i % 2; row = i // 2
    left = Inches(0.6) + col * Inches(6.3)
    top  = Inches(1.7) + row * Inches(2.7)
    add_panel(s, left, top, Inches(6.0), Inches(2.5))
    add_text(s, group, left + Inches(0.2), top + Inches(0.15),
             Inches(5.6), Inches(0.4), size=15, bold=True, color=ACCENT)
    add_text(s, items, left + Inches(0.3), top + Inches(0.6),
             Inches(5.6), Inches(1.9), size=12, color=INK, font="Consolas")

add_text(s, "Tools are pure Python — they mutate state or files but never call the LLM.",
         Inches(0.6), Inches(7.05), Inches(12), Inches(0.4),
         size=13, color=DIM, align=PP_ALIGN.CENTER)


# ---- 10. The LLM layer ---------------------------------------------------

s = new_slide()
slide_header(s, "Slide 9", "engine/llm.py — one provider switch, two backends")

add_bullets(s, [
    ("LLMConfig singleton",
     "Holds the active provider (ollama or gemini), model names, and the Ollama base URL. "
     "Seeded from .env, mutable at runtime."),
    ("call_llm(prompt, tools=None) — single entry point",
     "Everywhere else in the codebase imports just this function. "
     "Provider details are hidden behind it."),
    ("litellm under the hood",
     "Normalises the OpenAI / Gemini / Ollama tool-calling protocols into one shape, "
     "so the rest of the engine sees one schema."),
    ("Runtime toggle",
     "Press L in the Arcade window OR POST /api/llm/{provider} to flip Ollama ↔ Gemini "
     "without restarting. Useful when Ollama runs out of RAM mid-demo."),
    ("Default models",
     "Ollama: gemma4:e4b · Gemini: gemini-2.5-flash. "
     "We never call Anthropic — by design."),
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), size=14)


# ---- 11. Needs system ----------------------------------------------------

s = new_slide()
slide_header(s, "Slide 10", "engine/needs.py — the heartbeat of behaviour")

add_bullets(s, [
    ("Three numbers per agent",
     "hunger (0–100, higher = hungrier) · energy (0–100, lower = tireder) · mood (0–100)"),
    ("Decay over sim time",
     "Hunger +8/hour, energy −5/hour. Mood is event-driven — a refusal stings, "
     "a confirmed plan boosts."),
    ("Critical thresholds",
     "Hunger > 80 or energy < 20 → schedule is overridden, agent is told 'fix this now'. "
     "This is what stops them from working themselves to death."),
    ("Driven by sim time, not real time",
     "decay_needs(name, minutes_elapsed) is called from the tick loop with the "
     "delta in game-minutes — handles night auto-speed cleanly."),
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), size=14)


# ---- 12. Drama / story layer --------------------------------------------

s = new_slide()
slide_header(s, "Slide 11", "The story layer — drama you can watch unfold")

add_panel(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.4))

modules = [
    ("relationships.py",
     "Pairwise affinity scores. Updated when agents talk, refuse, agree, give items. "
     "Used by both the prompt and the viewer's relationship graph."),
    ("plots.py",
     "Tracks ongoing plot threads (e.g. 'Neha vs Vikram noise dispute'). "
     "Each thread has participants, status, and a list of beats."),
    ("narrator.py",
     "Periodic LLM call that reads recent events and produces a short, "
     "third-person commentary line — the in-game 'TV announcer'."),
    ("cliffhanger.py",
     "End-of-day hook: picks the most charged unresolved thread and writes a one-liner."),
    ("headlines.py",
     "Scans events for spicy interactions and converts them into gossip-style headlines."),
    ("protagonist.py",
     "Director-mode utility: nominates the day's lead character based on event volume."),
]
for i, (name, body) in enumerate(modules):
    col = i % 2; row = i // 2
    left = Inches(0.85) + col * Inches(5.95)
    top  = Inches(1.95) + row * Inches(1.65)
    add_text(s, name, left, top, Inches(5.7), Inches(0.4),
             size=15, bold=True, color=ACCENT, font="Consolas")
    add_text(s, body, left, top + Inches(0.45), Inches(5.7), Inches(1.2),
             size=12, color=INK)


# ---- 13. Time system ----------------------------------------------------

s = new_slide()
slide_header(s, "Slide 12", "Time — and why it matters")

# Big numbers row
def big(left, top, num, label):
    add_text(s, num, left, top, Inches(2.6), Inches(1.0),
             size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, label, left, top + Inches(0.95), Inches(2.6), Inches(0.5),
             size=13, color=DIM, align=PP_ALIGN.CENTER)

big(Inches(0.6),  Inches(1.7), "1 sec", "real second =\n5 game minutes")
big(Inches(3.6),  Inches(1.7), "3 sec", "= one tick =\n15 game minutes")
big(Inches(6.6),  Inches(1.7), "5 min", "real time =\n1 full game day")
big(Inches(9.6),  Inches(1.7), "96",   "ticks per\ngame day")

add_bullets(s, [
    ("Tick loop",
     "advance_time → decay_needs → asyncio.gather over all 10 agents → world.save() → sleep"),
    ("Per-agent schedules drive most decisions",
     "The injected SCHEDULE block tells the LLM 'right now, you sleep / commute / work / unwind'. "
     "Cuts down on bad decisions and saves tokens."),
    ("Drama-driven pacing",
     "compute_drama_score(state) weighs recent talk_to / conflict events, near-term plans, "
     "extreme moods, agents in motion. Low score → fast-forward to 4x. "
     "High score → snap back to 1x."),
    ("Night auto-speed",
     "When ≥7 agents are sleeping, the loop bumps speed to 4x; drops back when <5 sleep. "
     "The boring hours fly by."),
], Inches(0.6), Inches(3.5), Inches(12.1), Inches(3.5), size=13)


# ---- 14. World map ------------------------------------------------------

s = new_slide()
slide_header(s, "Slide 13", "world/map.json — eight connected locations")

locs = [
    ("apartment",     "Sushant Lok Apartments",  "home — sleep, rest"),
    ("metro",         "HUDA City Centre Metro",  "transit hub"),
    ("cyber_city",    "Cyber City",              "tech offices — work"),
    ("cyber_hub",     "Cyber Hub",               "boutiques, cafés"),
    ("sector29",      "Sector 29 Market",        "social, street food, earn"),
    ("dhaba",         "Pappu Dhaba",             "cheap meals, gossip"),
    ("supermarket",   "Star Bazaar",             "buy essentials"),
    ("park",          "Tau Devi Lal Park",       "exercise, calm"),
]
for i, (lid, name, role) in enumerate(locs):
    col = i % 4; row = i // 4
    left = Inches(0.6) + col * Inches(3.05)
    top  = Inches(1.7) + row * Inches(1.7)
    add_panel(s, left, top, Inches(2.95), Inches(1.55))
    add_text(s, lid, left + Inches(0.15), top + Inches(0.1),
             Inches(2.65), Inches(0.4), size=12, bold=True,
             color=ACCENT, font="Consolas")
    add_text(s, name, left + Inches(0.15), top + Inches(0.5),
             Inches(2.65), Inches(0.4), size=14, bold=True, color=INK)
    add_text(s, role, left + Inches(0.15), top + Inches(0.95),
             Inches(2.65), Inches(0.6), size=11, color=DIM)

add_text(s, "Each location has services (sleep, work, buy, earn, socialize…). "
            "move_to() runs BFS over the connection graph.",
         Inches(0.6), Inches(5.4), Inches(12), Inches(0.5),
         size=13, color=DIM, align=PP_ALIGN.CENTER)


# ---- 15. main.py renderer ------------------------------------------------

s = new_slide()
slide_header(s, "Slide 14", "main.py — the Arcade desktop window")

add_bullets(s, [
    ("Pure consumer of state",
     "Every frame: read WorldState, draw sprites, panels, dialog bubbles. "
     "Never writes back to state."),
    ("Background simulation thread",
     "asyncio.run(SimulationLoop().run()) on a daemon thread. "
     "Arcade's GL thread is left alone — important: Arcade is not async-friendly."),
    ("Sprite per agent",
     "Position interpolates between tick positions (tile_x, tile_y → pixels) "
     "for smooth movement."),
    ("Inspector panel",
     "Click an agent → side panel with diary, needs, last action, current goals."),
    ("Keyboard hotkeys",
     "Space pause · ←/→ slow/speed · F11 fullscreen · L toggle LLM · Esc close inspector."),
    ("Embedded web server",
     "main.py also boots server.py in another thread, so a browser viewer is always available."),
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), size=14)


# ---- 16. server.py + viewer.html -----------------------------------------

s = new_slide()
slide_header(s, "Slide 15", "server.py + viewer.html — the browser experience")

# Two columns
add_panel(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(5.4))
add_text(s, "server.py — FastAPI",
         Inches(0.8), Inches(1.85), Inches(5.6), Inches(0.5),
         size=18, bold=True, color=ACCENT2)
add_text(s, [
    "GET /api/state — full WorldState snapshot",
    "GET /api/agent/{name}/diary",
    "GET /api/agent/{name}/avatar",
    "GET /api/relationships",
    "GET /api/conversations?a=&b=",
    "POST /api/llm/{provider}  — switch backend",
    "GET / — serves viewer.html",
], Inches(0.8), Inches(2.4), Inches(5.6), Inches(4.5),
   size=13, color=INK, font="Consolas")
add_text(s, "Runs in a daemon uvicorn thread. Read-only proxy to WorldState — "
            "no simulation logic lives here.",
         Inches(0.8), Inches(6.2), Inches(5.6), Inches(0.8),
         size=12, color=DIM)

add_panel(s, Inches(6.8), Inches(1.7), Inches(6.0), Inches(5.4))
add_text(s, "viewer.html — vanilla JS",
         Inches(7.0), Inches(1.85), Inches(5.6), Inches(0.5),
         size=18, bold=True, color=GOOD)
add_text(s, [
    "HTML5 Canvas — no React, no build step",
    "Polls /api/state every ~1s",
    "Draws map, avatars, conversation log",
    "Sidebar: per-agent panel with diary",
    "Hover an avatar → enlarged portrait",
    "Click → inspector with goals + memory",
    "Tabs: Map · Relationships · Headlines",
], Inches(7.0), Inches(2.4), Inches(5.6), Inches(4.5),
   size=13, color=INK, font="Consolas")
add_text(s, "One file. ~180 lines. Loads from any device on the LAN — "
            "no client install.",
         Inches(7.0), Inches(6.2), Inches(5.6), Inches(0.8),
         size=12, color=DIM)


# ---- 17. Folder structure ------------------------------------------------

s = new_slide()
slide_header(s, "Slide 16", "Folder structure at a glance")

tree = """gurgaon_town_life/
├── main.py              # Arcade renderer + boots web server
├── server.py            # FastAPI app, daemon thread
├── viewer.html          # Browser viewer (vanilla JS / canvas)
│
├── engine/              # All simulation logic — no I/O frameworks
│   ├── world.py         # WorldState singleton, asyncio.Lock
│   ├── agent.py         # AgentRunner + LangGraph (4 nodes)
│   ├── tools.py         # All tool functions (move, talk, eat …)
│   ├── llm.py           # LLMConfig + call_llm — Ollama / Gemini
│   ├── needs.py         # decay_needs(); hunger/energy/mood
│   ├── relationships.py # Pairwise affinity tracking
│   ├── narrator.py      # Periodic commentary
│   ├── plots.py         # Plot-thread tracking
│   ├── cliffhanger.py   # End-of-day teaser
│   ├── headlines.py     # Gossip-style summaries
│   └── protagonist.py   # Director-mode lead picker
│
├── agents/              # One folder per character
│   ├── arjun/{soul,memory,diary,goals}.md
│   └── …                # × 10
│
├── world/
│   ├── map.json          # Locations + connections
│   ├── state.json        # Live snapshot — written each tick
│   └── scheduled_events.json
│
├── docs/                # tech_document.md, tech_stories.md
├── tests/               # Lightweight integrity tests
└── requirements.txt"""

add_code(s, tree, Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), size=12)


# ---- 18. State shape ----------------------------------------------------

s = new_slide()
slide_header(s, "Slide 17", "world/state.json — the persistent snapshot")

add_code(s,
'''{
  "day": 3, "sim_time": 540, "speed": 1,
  "agents": {
    "arjun": {
      "location": "cyber_city", "tile_x": 24, "tile_y": 15,
      "hunger": 42, "energy": 71, "mood": 55,
      "money": 1240, "inventory": ["bread", "chai"],
      "last_action": "working at cyber_city",
      "inbox": [{"from": "neha", "text": "lunch?"}]
    },
    "...": "..."
  },
  "events": [
    {"t": 540, "agent": "arjun", "text": "arjun says to neha: sure, 1pm"},
    {"t": 541, "agent": "neha",  "text": "neha confirms plan #7"}
  ],
  "shared_plans": [
    {"id": 7, "by": "neha", "with": "arjun",
     "where": "dhaba", "target_time": 4800, "status": "confirmed"}
  ],
  "relationships": {"arjun|neha": 0.62, "...": "..."}
}''',
Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), size=12)


# ---- 19. Concurrency model ----------------------------------------------

s = new_slide()
slide_header(s, "Slide 18", "Concurrency — three loops, one lock")

add_bullets(s, [
    ("Main thread (Arcade)",
     "Owns the OpenGL context. Draws frames at ~60 FPS by reading WorldState. "
     "Never blocks on I/O — never calls the LLM."),
    ("Simulation thread",
     "Daemon thread running its own asyncio loop. "
     "Each tick awaits asyncio.gather(*[runner.tick() for runner in agents]) — "
     "all 10 LLM calls fire in parallel."),
    ("Web server thread",
     "Daemon uvicorn thread. Read-only, returns JSON snapshots."),
    ("The lock",
     "WorldState._lock is an asyncio.Lock — only the simulation loop touches it. "
     "Renderers read with no lock; they may see a half-tick world but never a corrupt one, "
     "because the only writer batches its writes."),
    ("Atomic save",
     "world.save() writes to state.json.tmp then os.replace() — "
     "an interrupted save can never leave a partially written file."),
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), size=14)


# ---- 20. Design constraints --------------------------------------------

s = new_slide()
slide_header(s, "Slide 19", "Six golden rules baked into the codebase")

rules = [
    ("Engine ≠ renderer",       "Simulation logic NEVER imports arcade or fastapi."),
    ("Renderers are read-only", "The Arcade window and viewer.html never write state."),
    ("Tools never call the LLM","Tools mutate state — that's all. LLM calls happen in agent.py."),
    ("soul.md is immutable",    "Personality is git-versioned, not runtime-mutated."),
    ("All writes go through Lock", "Single asyncio.Lock — no race conditions."),
    ("No Anthropic API",        "By design — local-first via Ollama; Gemini optional."),
]
for i, (rule, body) in enumerate(rules):
    col = i % 2; row = i // 2
    left = Inches(0.6) + col * Inches(6.1)
    top  = Inches(1.7) + row * Inches(1.65)
    add_panel(s, left, top, Inches(5.95), Inches(1.5))
    add_text(s, rule, left + Inches(0.2), top + Inches(0.15),
             Inches(5.6), Inches(0.5), size=15, bold=True, color=ACCENT)
    add_text(s, body, left + Inches(0.2), top + Inches(0.65),
             Inches(5.6), Inches(0.85), size=13, color=INK)


# ---- 21. Tech stack summary ---------------------------------------------

s = new_slide()
slide_header(s, "Slide 20", "Tech stack — and why each piece was chosen")

stack = [
    ("Python 3.12",   "asyncio + match statements + good typing — single language end-to-end.",        ACCENT),
    ("LangGraph",     "Typed state graph for the per-agent decision loop. Easy to extend.",           ACCENT2),
    ("LiteLLM",       "Single tool-calling shape across Ollama / Gemini / OpenAI.",                   GOOD),
    ("Ollama",        "Local LLM (gemma4:e4b) — free, private, no rate limits.",                       RGBColor(0xE6, 0x6E, 0xC9)),
    ("Gemini 2.5 Flash","Cloud fallback — kicks in when Ollama is too slow or out of RAM.",            ACCENT),
    ("Arcade 3",      "Pythonic 2D game engine — sprites, camera, easy keyboard handling.",           ACCENT2),
    ("FastAPI + uvicorn","Web API layer — async, OpenAPI docs for free, runs in a thread.",            GOOD),
    ("Vanilla JS / Canvas","Zero-build viewer. Loads on any phone on the LAN.",                        RGBColor(0xE6, 0x6E, 0xC9)),
]
for i, (name, why, color) in enumerate(stack):
    col = i % 2; row = i // 2
    left = Inches(0.6) + col * Inches(6.1)
    top  = Inches(1.7) + row * Inches(1.30)
    add_panel(s, left, top, Inches(5.95), Inches(1.18))
    add_chip(s, name, left + Inches(0.2), top + Inches(0.18),
             Inches(2.0), Inches(0.4), fill=color, color=BG, size=11)
    add_text(s, why, left + Inches(0.2), top + Inches(0.65),
             Inches(5.6), Inches(0.6), size=12, color=INK)


# ---- 22. What's clever / talking points --------------------------------

s = new_slide()
slide_header(s, "Slide 21", "What makes this codebase interesting")

add_bullets(s, [
    ("Files-as-mind",
     "An agent's identity is four markdown files. You can git-diff a personality. "
     "You can edit Arjun's goals.md in VS Code and watch him pivot in real time."),
    ("LLM context is dynamic but bounded",
     "We don't dump all of memory.md into every prompt — we grep_memory() for keywords "
     "from the current scene. Prompts stay small even on day 30."),
    ("Schedule pre-prompt cuts hallucinations",
     "Hard-coding 'right now you're at work' beats hoping the LLM infers it from the clock. "
     "Saves tokens AND produces saner behaviour."),
    ("Drama-driven pacing",
     "Score the world's recent activity. Quiet world → fast-forward 4x. Spicy world → 1x. "
     "The viewer is never bored."),
    ("Zero-build web viewer",
     "viewer.html is one file, ~180 lines, no npm. Anyone on the LAN can watch."),
    ("Three layers, one lock",
     "Engine, renderer, web server share state through a single asyncio-locked dict. "
     "No queues, no events, no observers — just one snapshot they all read."),
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), size=13)


# ---- 23. Limitations / future ------------------------------------------

s = new_slide()
slide_header(s, "Slide 22", "Honest limitations & where it could go next")

add_bullets(s, [
    ("Ollama is single-machine",
     "All 10 agents share the same model + VRAM. On a laptop GPU this caps the tick rate."),
    ("No long-term planning across days",
     "Goals reset each morning. Agents don't really hold multi-day arcs without help from plots.py."),
    ("Mood is a single number",
     "It works, but a vector (anger / joy / anxiety / boredom) would give richer behaviour."),
    ("Memory is unbounded text",
     "memory.md grows over time. No summarisation pass yet — eventually the file gets fat."),
    ("Renderer is laptop-grade",
     "Arcade is fine for 10 sprites; it would not scale to 200 agents without batching."),
    ("Nothing tested past Day 5",
     "Long-run behaviours (relationship drift, stale goals) need a soak test we haven't run."),
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), size=14)


# ===========================================================================
# INTERVIEW SECTION
# ===========================================================================

# ---- 24. Interview section divider --------------------------------------

s = new_slide()
add_bg(s, RGBColor(0x08, 0x0C, 0x14))
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.4),
                         prs.slide_width, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background(); bar.shadow.inherit = False

add_text(s, "PART 2", Inches(0.8), Inches(2.0), Inches(12), Inches(0.5),
         size=14, bold=True, color=ACCENT)
add_text(s, "How to present this in an interview",
         Inches(0.8), Inches(2.4), Inches(12), Inches(1.0),
         size=44, bold=True, color=INK)
add_text(s, "Story arc · 30-second pitch · likely questions · demo flow · skill-by-skill talking points",
         Inches(0.8), Inches(3.7), Inches(12), Inches(0.5),
         size=18, color=DIM)


# ---- 25. The 30-second pitch --------------------------------------------

s = new_slide()
slide_header(s, "Interview · 1", "The 30-second pitch (memorise this)")

add_panel(s, Inches(1.0), Inches(1.9), Inches(11.3), Inches(3.8))
add_text(s,
         "“It's a multi-agent simulation where ten LLM-powered characters live "
         "their own lives in a small virtual town. Each one has a personality file, "
         "a private diary, and goals they rewrite each morning. Every three seconds "
         "all ten run a four-node decision graph in parallel — gather context, "
         "ask the LLM, execute one tool, reflect. Their state lives in a single "
         "JSON file behind an asyncio lock, and you can watch them either in an "
         "Arcade desktop window or in a browser viewer. The interesting parts are "
         "the file-based 'mind' design, the per-archetype schedule injection that "
         "keeps prompts short, and the drama-score pacer that fast-forwards the "
         "boring stretches.”",
         Inches(1.3), Inches(2.1), Inches(10.7), Inches(3.4),
         size=16, color=INK)

add_text(s, "Read it twice. Three sentences if cut short: WHAT it is, HOW it runs, WHY it's clever.",
         Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.5),
         size=14, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)


# ---- 26. Story arc -------------------------------------------------------

s = new_slide()
slide_header(s, "Interview · 2", "Tell it as a story — the 4-act arc")

acts = [
    ("ACT 1 — The Problem",
     "“I wanted to see what would happen if you just let LLMs live a life — "
     "no tasks, no users, just personality plus a clock. Agents in the wild are "
     "usually goal-directed assistants. I wanted observers, not operators.”"),
    ("ACT 2 — The Architecture",
     "“So I split it into three layers: an engine that owns all state, a renderer "
     "that only reads, and a web server that's basically a mirror. The engine is "
     "where everything interesting happens — a single WorldState dict behind an "
     "asyncio lock, with a four-node LangGraph per agent.”"),
    ("ACT 3 — The Tricky Bits",
     "“Three things were hard: keeping prompts small as memory grows (solved with "
     "grep_memory), keeping behaviour sane around the clock (solved with archetype-"
     "based schedule injection), and keeping the viewer interesting (solved with a "
     "drama-score auto-pacer).”"),
    ("ACT 4 — What I'd Do Next",
     "“The next step is multi-day arcs. Right now plot threads exist but agents "
     "don't really plan past tomorrow. I'd add a weekly-planning node and a "
     "memory-summarisation pass so memory.md stays bounded.”"),
]
for i, (title, body) in enumerate(acts):
    col = i % 2; row = i // 2
    left = Inches(0.6) + col * Inches(6.1)
    top  = Inches(1.7) + row * Inches(2.65)
    add_panel(s, left, top, Inches(5.95), Inches(2.5))
    add_text(s, title, left + Inches(0.2), top + Inches(0.15),
             Inches(5.6), Inches(0.4), size=14, bold=True, color=ACCENT)
    add_text(s, body, left + Inches(0.2), top + Inches(0.65),
             Inches(5.6), Inches(1.85), size=12, color=INK)


# ---- 27. Demo flow -------------------------------------------------------

s = new_slide()
slide_header(s, "Interview · 3", "If you get to do a live demo — the script")

add_bullets(s, [
    ("0:00 — Open the browser viewer at localhost:8000",
     "Avoid Arcade if screen-sharing — the GL window can flicker over Zoom. The web viewer is bullet-proof."),
    ("0:30 — Point at the map",
     "Name 3 characters. 'Arjun is at Cyber City — that's his office. Suresh is at the metro — he drives a cab.' Establish that these are not scripted."),
    ("1:00 — Open Arjun's diary",
     "Show that this file is being written by the LLM, not by you. Read one entry aloud."),
    ("1:30 — Click an avatar → inspector",
     "Show needs, last action, current goals. Note: 'goals were rewritten this morning, by him.'"),
    ("2:00 — Show the headlines tab / relationship graph",
     "Drama you didn't author. This is the wow moment."),
    ("2:30 — Open VS Code, edit Arjun's goals.md",
     "Add 'pick a fight with Vikram today'. Save. Wait one tick. Watch his next action shift. This is the killer demo."),
    ("3:00 — Wrap with the architecture slide",
     "Three layers, one lock, file-based mind. Ask if they want to dive into any one of them."),
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), size=13)


# ---- 28. Likely questions -----------------------------------------------

s = new_slide()
slide_header(s, "Interview · 4", "Questions you WILL be asked — with answers")

qs = [
    ("Why LangGraph instead of just calling the LLM in a loop?",
     "Because the loop is a state machine: gather → decide → execute → reflect. "
     "LangGraph gives me typed state and named nodes, so adding (e.g.) a planning node "
     "doesn't ripple through. It's overkill for one node, just right for four."),
    ("Why local Ollama and not OpenAI / Anthropic?",
     "Cost + privacy + offline. Ten parallel calls per tick on a paid API would "
     "bankrupt the demo. Ollama is free and the gemma4:e4b model is good enough "
     "for short, structured tool-calls. Gemini is the cloud fallback."),
    ("How do you keep the LLM in character?",
     "Three things: (1) soul.md is injected verbatim every prompt — short and dense; "
     "(2) we grep memory.md for keywords from the current scene rather than dumping it; "
     "(3) the schedule pre-prompt anchors them in time-of-day, so they don't go to work at 3am."),
    ("How do ten agents not corrupt the state file?",
     "Single asyncio.Lock around every WorldState mutation. Reads are unlocked but the "
     "writer batches its updates per tick, so a reader at worst sees a half-written tick — "
     "never a malformed JSON file. world.save() uses tmp + os.replace() so the file on "
     "disk is always atomic."),
    ("What's the hardest bug you fixed?",
     "Pick a real one — e.g. agents loop-talking to themselves at 3am because "
     "the night auto-speed advanced sim_time faster than mood could decay. "
     "Fixed by gating critical-need overrides behind sleep windows."),
    ("How would you scale this to 100 agents?",
     "Three changes: batch the LLM calls (one prompt with N agents), shard WorldState by "
     "neighbourhood with a lock per shard, and replace the JSON snapshot with SQLite WAL "
     "for the event log."),
]
# Two columns
for i, (q, a) in enumerate(qs):
    col = i % 2; row = i // 2
    left = Inches(0.6) + col * Inches(6.1)
    top  = Inches(1.7) + row * Inches(1.85)
    add_panel(s, left, top, Inches(5.95), Inches(1.7))
    add_text(s, "Q. " + q, left + Inches(0.2), top + Inches(0.15),
             Inches(5.6), Inches(0.55), size=12, bold=True, color=ACCENT)
    add_text(s, a, left + Inches(0.2), top + Inches(0.65),
             Inches(5.6), Inches(1.05), size=11, color=INK)


# ---- 29. Skill-by-skill talking points -----------------------------------

s = new_slide()
slide_header(s, "Interview · 5", "Map the project to the skills they're hiring for")

skills = [
    ("Concurrency",
     "asyncio.Lock around WorldState · asyncio.gather over 10 agents per tick · "
     "daemon thread for the web server · atomic save with os.replace()"),
    ("LLM / agentic",
     "LangGraph state machine · tool-call schemas via litellm · prompt assembly "
     "with grep-style memory retrieval · runtime provider switching"),
    ("System design",
     "Three-layer split (engine / renderer / web) · single source of truth · "
     "files-as-mind for explainability · drama-score auto-pacing"),
    ("API design",
     "FastAPI with read-only resource endpoints · POST /api/llm/{provider} for "
     "runtime control · zero-build vanilla-JS client polling /api/state"),
    ("Game / sim engineering",
     "BFS pathfinding for move_to · per-archetype schedules · sprite interpolation "
     "between tick positions · click-to-inspect UI"),
    ("Operability",
     "All state in one human-readable JSON · agent minds in markdown — git-diffable · "
     "--reset flag · runtime LLM toggle for live recovery"),
]
for i, (name, body) in enumerate(skills):
    col = i % 2; row = i // 3
    left = Inches(0.6) + col * Inches(6.1)
    top  = Inches(1.7) + row * Inches(1.75)
    add_panel(s, left, top, Inches(5.95), Inches(1.6))
    add_text(s, name, left + Inches(0.2), top + Inches(0.15),
             Inches(5.6), Inches(0.4), size=14, bold=True, color=ACCENT)
    add_text(s, body, left + Inches(0.2), top + Inches(0.6),
             Inches(5.6), Inches(0.95), size=12, color=INK)


# ---- 30. Pitfalls to avoid ----------------------------------------------

s = new_slide()
slide_header(s, "Interview · 6", "Pitfalls — what NOT to say")

add_bullets(s, [
    ("Don't oversell it as 'AGI' or 'emergent consciousness'",
     "It's a fun simulation built on small LLM calls. Engineers smell hype — "
     "describe it as a system, not a miracle."),
    ("Don't claim 100% reliability",
     "Be honest: at 3am with a hot Ollama, agents do hallucinate. Acknowledging "
     "limits scores higher than pretending there are none."),
    ("Don't say 'it's just a prompt'",
     "Undersells the engineering. The interesting work is around the prompt — "
     "schedules, tools, memory grep, drama pacing, locking, persistence."),
    ("Don't get lost in personalities",
     "Recruiters might ooh-and-aah over Neha and Vikram. Smile, then steer back "
     "to architecture in 15 seconds."),
    ("Don't promise features that aren't built",
     "If they ask 'does it remember last week?' — say 'no, that's the next thing I'd add', "
     "not 'kind of, sort of'. Honest beats vague."),
    ("Don't show the Arcade window over Zoom",
     "GL surfaces flicker. The browser viewer is what you demo."),
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), size=14)


# ---- 31. One-slide cheat sheet -------------------------------------------

s = new_slide()
slide_header(s, "Interview · 7", "One-slide cheat sheet — print this")

add_panel(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5))

add_text(s, "WHAT", Inches(0.9), Inches(1.85), Inches(2.0), Inches(0.4),
         size=14, bold=True, color=ACCENT)
add_text(s, "10 LLM agents living in a virtual town. No player. Pure observation.",
         Inches(2.9), Inches(1.85), Inches(9.6), Inches(0.5), size=14, color=INK)

add_text(s, "STACK", Inches(0.9), Inches(2.35), Inches(2.0), Inches(0.4),
         size=14, bold=True, color=ACCENT)
add_text(s, "Python 3.12 · LangGraph · asyncio · FastAPI · Arcade · LiteLLM (Ollama / Gemini)",
         Inches(2.9), Inches(2.35), Inches(9.6), Inches(0.5), size=14, color=INK)

add_text(s, "LOOP", Inches(0.9), Inches(2.85), Inches(2.0), Inches(0.4),
         size=14, bold=True, color=ACCENT)
add_text(s, "Per agent, every 3s: gather_context → llm_decide → execute_tool → reflect",
         Inches(2.9), Inches(2.85), Inches(9.6), Inches(0.5), size=14, color=INK)

add_text(s, "MIND", Inches(0.9), Inches(3.35), Inches(2.0), Inches(0.4),
         size=14, bold=True, color=ACCENT)
add_text(s, "soul.md (immutable) + memory.md + diary.md + goals.md — all markdown, all editable",
         Inches(2.9), Inches(3.35), Inches(9.6), Inches(0.5), size=14, color=INK)

add_text(s, "STATE", Inches(0.9), Inches(3.85), Inches(2.0), Inches(0.4),
         size=14, bold=True, color=ACCENT)
add_text(s, "world/state.json — one dict, asyncio.Lock, atomic saves",
         Inches(2.9), Inches(3.85), Inches(9.6), Inches(0.5), size=14, color=INK)

add_text(s, "VIEWS", Inches(0.9), Inches(4.35), Inches(2.0), Inches(0.4),
         size=14, bold=True, color=ACCENT)
add_text(s, "Arcade desktop window + browser viewer — both read-only consumers",
         Inches(2.9), Inches(4.35), Inches(9.6), Inches(0.5), size=14, color=INK)

add_text(s, "CLEVER", Inches(0.9), Inches(4.85), Inches(2.0), Inches(0.4),
         size=14, bold=True, color=ACCENT)
add_text(s, "Files-as-mind · schedule pre-prompt · grep_memory · drama-score auto-pacing",
         Inches(2.9), Inches(4.85), Inches(9.6), Inches(0.5), size=14, color=INK)

add_text(s, "DEMO",  Inches(0.9), Inches(5.35), Inches(2.0), Inches(0.4),
         size=14, bold=True, color=ACCENT)
add_text(s, "Browser → click avatar → diary → edit goals.md live → watch behaviour shift",
         Inches(2.9), Inches(5.35), Inches(9.6), Inches(0.5), size=14, color=INK)

add_text(s, "WEAK",  Inches(0.9), Inches(5.85), Inches(2.0), Inches(0.4),
         size=14, bold=True, color=ACCENT)
add_text(s, "No multi-day planning · unbounded memory.md · single-machine Ollama",
         Inches(2.9), Inches(5.85), Inches(9.6), Inches(0.5), size=14, color=INK)

add_text(s, "NEXT",  Inches(0.9), Inches(6.35), Inches(2.0), Inches(0.4),
         size=14, bold=True, color=ACCENT)
add_text(s, "Weekly-planning node · memory summariser · 100-agent batching",
         Inches(2.9), Inches(6.35), Inches(9.6), Inches(0.5), size=14, color=INK)


# ---- 32. Closing slide ---------------------------------------------------

s = new_slide()
add_bg(s, RGBColor(0x08, 0x0C, 0x14))
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.4),
                         prs.slide_width, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background(); bar.shadow.inherit = False

add_text(s, "Thanks for watching the town", Inches(0.8), Inches(2.2),
         Inches(12), Inches(1.0), size=42, bold=True, color=INK)
add_text(s, "Questions? Open agents/<name>/diary.md and read along.",
         Inches(0.8), Inches(3.7), Inches(12), Inches(0.5),
         size=20, color=DIM)
add_text(s, "github · docs/tech_document.md · CLAUDE.md",
         Inches(0.8), Inches(4.4), Inches(12), Inches(0.5),
         size=14, color=ACCENT)


# ---- Footers (skip title + section divider + closing) -------------------

total = len(slides_built)
for i, sl in enumerate(slides_built, start=1):
    if i in (1, total):
        continue
    add_footer(sl, i, total)


# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------

out = pathlib.Path(__file__).resolve().parent.parent / "docs" / "gurgaon_town_life_explained.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"Wrote {out}  ({total} slides)")
